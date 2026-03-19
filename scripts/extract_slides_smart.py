#!/usr/bin/env python3
"""
Smart PPT slide extractor - extracts images per slide
Pure Python, no external dependencies except python-pptx
"""

import sys
from pathlib import Path
from pptx import Presentation
from PIL import Image
import io


def extract_slide_images(ppt_path, output_dir):
    """
    Extract images from each slide, organized by slide number
    If a slide has one main image, extract it as that slide's screenshot
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(ppt_path)

    print(f"PPT has {len(prs.slides)} slides\n")

    slide_images = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        print(f"Processing slide {slide_idx}...")

        # Find all images in this slide
        images_in_slide = []

        for shape in slide.shapes:
            if hasattr(shape, "image"):
                # This shape contains an image
                image = shape.image
                image_bytes = image.blob

                # Get image dimensions to find the largest one
                img = Image.open(io.BytesIO(image_bytes))
                width, height = img.size

                images_in_slide.append({
                    'image': img,
                    'bytes': image_bytes,
                    'size': width * height,
                    'ext': image.ext
                })

        if images_in_slide:
            # Find the largest image (likely the main screenshot)
            largest = max(images_in_slide, key=lambda x: x['size'])

            # Save as slide_NNN.png
            output_path = output_dir / f"slide_{slide_idx:03d}.png"
            largest['image'].save(output_path, 'PNG')

            print(f"  ✓ Saved: slide_{slide_idx:03d}.png ({largest['image'].size[0]}x{largest['image'].size[1]})")
            slide_images.append(str(output_path))
        else:
            print(f"  ⚠ No images found in slide {slide_idx}")

    return slide_images


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python3 extract_slides_smart.py <pptx_file> <output_dir>")
        sys.exit(1)

    ppt_file = sys.argv[1]
    output_dir = sys.argv[2]

    print(f"Input:  {ppt_file}")
    print(f"Output: {output_dir}\n")

    images = extract_slide_images(ppt_file, output_dir)

    print(f"\n✓ Extracted {len(images)} slide images")
    print(f"Location: {output_dir}")
