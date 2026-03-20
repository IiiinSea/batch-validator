#!/usr/bin/env python3
"""
Extract text content from each PPT slide
"""

import sys
from pathlib import Path
from pptx import Presentation
import json


def extract_text_from_slide(slide):
    """Extract all text from a slide"""
    text_content = []

    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text_content.append(shape.text.strip())

    return text_content


def extract_ppt_text(ppt_path, output_json=None):
    """
    Extract text from all slides in PPT

    Args:
        ppt_path: Path to the .pptx file
        output_json: Optional path to save JSON output

    Returns:
        Dictionary with slide text content
    """
    prs = Presentation(ppt_path)

    slides_data = {}

    for slide_idx, slide in enumerate(prs.slides, start=1):
        text_content = extract_text_from_slide(slide)

        slides_data[f"slide_{slide_idx:03d}"] = {
            "slide_number": slide_idx,
            "text_lines": text_content,
            "full_text": "\n".join(text_content)
        }

        print(f"Slide {slide_idx}:")
        for line in text_content:
            print(f"  - {line}")
        print()

    # Save to JSON if requested
    if output_json:
        output_path = Path(output_json)
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(slides_data, f, ensure_ascii=False, indent=2)
        print(f"\n✓ Text content saved to: {output_path}")

    return slides_data


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="从 PPT 中提取每页文字")
    parser.add_argument("pptx", help="PPT 文件路径")
    parser.add_argument("--json", dest="output_json", help="输出 JSON 文件路径")
    parser.add_argument(
        "--txt-dir",
        help="每页输出一个 txt 文件的目录（不指定则自动用 {PPT名}_text/）",
    )
    args = parser.parse_args()

    print(f"Extracting text from: {args.pptx}\n")
    print("=" * 80)

    slides_data = extract_ppt_text(args.pptx, args.output_json)

    print("=" * 80)
    print(f"\n✓ Extracted text from {len(slides_data)} slides")

    txt_dir = args.txt_dir or str(Path(args.pptx).parent / f"{Path(args.pptx).stem}_text")
    out = Path(txt_dir)
    out.mkdir(parents=True, exist_ok=True)
    for key, data in slides_data.items():
        (out / f"{key}.txt").write_text(
            f"=== Slide {data['slide_number']} ===\n{data['full_text']}\n",
            encoding="utf-8",
        )
    print(f"✓ Per-slide txt files saved to: {out}")
