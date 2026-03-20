#!/usr/bin/env python3
"""
从 PPT 中提取所有图片，每页一个文件夹。

用法:
  uv run scripts/extract_all_images.py <pptx_file> [output_dir]
"""

import sys
from pathlib import Path
from pptx import Presentation
from PIL import Image
import io


def extract_all_images(ppt_path: str, output_dir: str | None = None) -> dict[int, list[str]]:
    ppt = Path(ppt_path)
    out = Path(output_dir) if output_dir else ppt.parent / f"{ppt.stem}_images"
    out.mkdir(parents=True, exist_ok=True)

    prs = Presentation(ppt_path)
    print(f"PPT 共 {len(prs.slides)} 页\n")

    result: dict[int, list[str]] = {}

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_dir = out / f"slide_{slide_idx:03d}"
        slide_dir.mkdir(exist_ok=True)

        images_saved = []
        img_seq = 0

        for shape in slide.shapes:
            if not hasattr(shape, "image"):
                continue

            img_seq += 1
            img = Image.open(io.BytesIO(shape.image.blob))
            w, h = img.size
            filename = f"img_{img_seq:02d}_{w}x{h}.png"
            save_path = slide_dir / filename
            img.save(save_path, "PNG")
            images_saved.append(str(save_path))

        result[slide_idx] = images_saved

        if images_saved:
            print(f"  第 {slide_idx} 页: {len(images_saved)} 张图片 → {slide_dir.name}/")
        else:
            print(f"  第 {slide_idx} 页: 无图片")
            slide_dir.rmdir()

    total = sum(len(v) for v in result.values())
    print(f"\n✓ 共提取 {total} 张图片 → {out}")
    return result


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python extract_all_images.py <pptx_file> [output_dir]")
        sys.exit(1)

    extract_all_images(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
