#!/usr/bin/env python3
"""
Extract ALL images from each PPT slide (not just the largest one)
Each slide may have multiple screenshots containing different information
"""

import sys
from pathlib import Path
from pptx import Presentation
from PIL import Image
import io


def extract_all_images_per_slide(ppt_path, output_dir):
    """
    Extract ALL images from each slide, organized by slide number

    Args:
        ppt_path: Path to the .pptx file
        output_dir: Directory to save the slide images

    Returns:
        Dict mapping slide number to list of image paths
        Example: {1: ['slide_001_img_01.png', 'slide_001_img_02.png'], ...}
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(ppt_path)

    print(f"PPT has {len(prs.slides)} slides\n")

    slides_images_map = {}

    for slide_idx, slide in enumerate(prs.slides, start=1):
        print(f"Processing slide {slide_idx}...")

        # Find all images in this slide
        images_in_slide = []
        image_count = 0

        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob

                # Get image
                img = Image.open(io.BytesIO(image_bytes))
                width, height = img.size

                image_count += 1

                # Save as slide_NNN_img_MM.png
                output_path = output_dir / f"slide_{slide_idx:03d}_img_{image_count:02d}.png"
                img.save(output_path, 'PNG')

                images_in_slide.append(str(output_path))

                print(f"  ✓ Image {image_count}: {output_path.name} ({width}x{height})")

        if not images_in_slide:
            print(f"  ⚠ No images found in slide {slide_idx}")

        slides_images_map[slide_idx] = images_in_slide

    return slides_images_map


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python3 extract_all_slide_images.py <pptx_file> <output_dir>")
        print()
        print("Extracts ALL images from each slide (not just the largest one)")
        sys.exit(1)

    ppt_file = sys.argv[1]
    output_dir = sys.argv[2]

    print(f"Input:  {ppt_file}")
    print(f"Output: {output_dir}\n")

    slides_map = extract_all_images_per_slide(ppt_file, output_dir)

    print("\n" + "=" * 80)
    print("Summary:")
    total_images = sum(len(imgs) for imgs in slides_map.values())
    print(f"  Total slides: {len(slides_map)}")
    print(f"  Total images: {total_images}")
    print(f"  Average images per slide: {total_images / len(slides_map):.1f}")

    # Save mapping to JSON
    import json
    mapping_file = Path(output_dir) / "slides_images_map.json"

    # Convert integer keys to strings for JSON
    json_map = {str(k): v for k, v in slides_map.items()}

    with open(mapping_file, 'w', encoding='utf-8') as f:
        json.dump(json_map, f, indent=2, ensure_ascii=False)

    print(f"\n✓ Image mapping saved to: {mapping_file}")
